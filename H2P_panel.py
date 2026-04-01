"""
To define the main window of H2P as a class
This file also contains all the (callback) routines
that
    write_presentation() .. by
    open_presentation() .. open an existing base pptx
    for each hymn
        for each verse or chorus
            inject_verse (or choruses) into  a new textbox on a new slide
            inject hymn_number() in a top right-hand box 
        inject_attribution() at end of each hymn, inject_attribution()
        inject_blank_slide() between verses
    
CJCJ 13-03-2026
"""

import sys

import webbrowser
import os
import time

from pathlib import Path


from PySide6.QtGui import (
    QAction,
    QFont,
    QFontDatabase,
    QIntValidator,
    QRegularExpressionValidator,
    QIcon)

from PySide6.QtWidgets import (
    QApplication,
    QHBoxLayout,
    QMainWindow,
    QVBoxLayout,
    QGridLayout,
    QWidget,
    QComboBox,
    QTextEdit,
    QLineEdit,
    QPushButton,
    QLabel,
    QFileDialog,
    QFrame, 
    QMessageBox,
    QDialog,
    QCompleter)

from PySide6.QtCore import (
    QRect,
    QSize,
    Signal,QRegularExpression)

# some constants for formatting
MIN_NUMBER_WIDTH = 24
EXACT_NUMBER_WIDTH = 90
MAX_HYMN_NUMBER_WIDTH = 60
MIN_LINE_WIDTH = 460
FONT_NAME = "Andika" # for interface
INTERFACE_FONT_SIZE = 13

# import so that cfg 'globals' are available in all functions here.

import H2P_cfg as cfg

class MainWindow(QMainWindow):

    def __init__(self):
        super().__init__()

        # set up a nice font for the interface (The presentation may use a different font)
        fontPath = str(Path('.','H2P-private-data','Andika-Regular.ttf'))
        id = QFontDatabase.addApplicationFont(fontPath)
        if id < 0:
            print("Error reading font, Andika-Regular.ttf")
            cfg.families = QFontDatabase.families()
        else:
            cfg.families = QFontDatabase.applicationFontFamilies(id)

            # apply font and sizee to window title 
            self.setWindowTitle("Hymns \u2192 Presentation")
            self.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE)) # this sets text font for all children of self, the MainWindow (but not the window title - because that belongs to the desktop(?))
            iconPath = str(Path('.','H2P-private-data','H2P-ico.png'))
            self.setWindowIcon(QIcon(iconPath))
        
        # load graphic assets to buttons

        self.savePptxButton = QPushButton()
        self.savePptxButton.setFlat(True)

        # set up the menus
        
        menu = self.menuBar()
        menu.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE)) # sets all menu items
        file_menu = menu.addMenu("&File")
        file_menu.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        
        menuActionReload = QAction("Reload hymn database", self)
        menuActionReload.setStatusTip("Poll hymn txt files in subdirectories of hymns.")
        menuActionReload.triggered.connect(self.Reload)
        menuActionReload.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        file_menu.addAction(menuActionReload)
        
        menuActionIndex = QAction("Create hymn index", self)
        menuActionIndex.setStatusTip("Show an index of first lines for a source books.")
        menuActionIndex.triggered.connect(self.MakeIndex)
        menuActionIndex.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        file_menu.addAction(menuActionIndex)

        menuActionLog = QAction("Show log file", self)
        menuActionLog.setStatusTip("CPresent the hymn-reading log in the webbrowser.")
        menuActionLog.triggered.connect(self.showLog)
        menuActionLog.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        file_menu.addAction(menuActionLog)

        menuActionExit = QAction("&Exit", self)
        menuActionExit.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        menuActionExit.setStatusTip("Exit the program")
        menuActionExit.triggered.connect(self.Exit)
        file_menu.addAction(menuActionExit)
        
        settings_menu = menu.addMenu("&Settings")
        settings_menu.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        menuActionSet = QAction("CCLI", self)
        menuActionSet.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        menuActionSet.setStatusTip("Set licence No. to be displayed in presentation.")
        menuActionSet.triggered.connect(self.setCCLI)
        settings_menu.addAction(menuActionSet)

        self.menuActionBlack = QAction("Black divider", self)
        self.menuActionBlack.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        self.menuActionBlack.setStatusTip("Slide between hymns black or same background as verse slides.")
        self.menuActionBlack.setCheckable(True)
        self.menuActionBlack.setChecked(cfg.black_between)
        self.menuActionBlack.setEnabled(True)
        self.menuActionBlack.toggled.connect(self.recordBlackBetween)
        settings_menu.addAction(self.menuActionBlack)

        about_menu = menu.addMenu("&Help")
        menuActionHelp = QAction("Help", self)
        menuActionHelp.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        menuActionHelp.setStatusTip("Open the software manual.")
        menuActionHelp.triggered.connect(self.Help)
        about_menu.addAction(menuActionHelp)
        
        menuActionAbout = QAction("About", self)
        menuActionAbout.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        menuActionAbout.setStatusTip("Show author details.")
        menuActionAbout.triggered.connect(self.About)
        about_menu.addAction(menuActionAbout)
        
        # set up layouts for left, right
        # hymn selction layouts

                
        layout1 = QHBoxLayout()
        layout1.setContentsMargins(4, 4, 4, 4)
        layout1.setSpacing(4)

        layout2 = QHBoxLayout()
        layout2.setContentsMargins(4, 4, 4, 4)
        layout2.setSpacing(4)
        
        layout3 = QHBoxLayout()
        layout3.setContentsMargins(4, 4, 4, 4)
        layout3.setSpacing(4)
        
        layout4 = QHBoxLayout()
        layout4.setContentsMargins(4, 4, 4, 4)
        layout4.setSpacing(4)
        
        layout5 = QHBoxLayout()
        layout5.setContentsMargins(4, 4, 4, 4)
        layout5.setSpacing(4)
        
        layout6 = QHBoxLayout()
        layout6.setContentsMargins(4, 4, 4, 4)
        layout6.setSpacing(4)
        
        layout7 = QHBoxLayout()
        layout7.setContentsMargins(4, 4, 4, 4)
        layout7.setSpacing(4)

        layoutLeft = QVBoxLayout()
        layoutLeft.setContentsMargins(4, 4, 4, 4)
        layoutLeft.setSpacing(4)
 
        # the text fields for left side
        
        labelLeftHeader = QLabel("Choose hymns")
        labelLeftHeader.setStyleSheet(cfg.colStr) # changing colour resets font size
        labelLeftHeader.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE)) # so have to change it back

        layoutLeft.addWidget(labelLeftHeader)
        
        # Hymn selction layouts

        # hymn number validator only has to be set up once and then applied to each hymn number QLineEdit.
        RegX = QRegularExpression("[1-9]\\d{0,3}|") # yea it seems to work - accepts up to 4 digits or a blank string
        validNumber = QRegularExpressionValidator(RegX) #Works better than QIntValidator which does not allow a blank.
        
        numberLabel1 = QLabel("1.")
        numberLabel1.setStyleSheet(cfg.colStr)
        numberLabel1.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel1.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn1 = QComboBox()
        self.comboHymn1.setEditable(False)
        self.comboHymn1.addItems(cfg.bookCodeList)
        self.comboHymn1.currentIndexChanged.connect(self.BookUpdate1)

        self.lineEditHymn1 = QLineEdit("0")
        self.lineEditHymn1.setValidator(validNumber) 
        self.lineEditHymn1.setClearButtonEnabled(True)
        self.lineEditHymn1.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn1.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn1.editingFinished.connect(self.UpdateHymnNumber1)

        self.textFirstLine1 = QLineEdit()
        self.textFirstLine1.setText("Set hymn book and number to check first line.")
        self.textFirstLine1.setReadOnly(True)
        self.textFirstLine1.setMinimumWidth(MIN_LINE_WIDTH)
        
        layout1.addWidget(numberLabel1)
        layout1.addWidget(self.comboHymn1)
        layout1.addWidget(self.lineEditHymn1)
        layout1.addWidget(self.textFirstLine1)
        layoutLeft.addLayout(layout1)

        numberLabel2 = QLabel("2.")
        numberLabel2.setStyleSheet(cfg.colStr)
        numberLabel2.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel2.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn2 = QComboBox()
        self.comboHymn2.setEditable(False)
        self.comboHymn2.addItems(cfg.bookCodeList)

        self.lineEditHymn2 = QLineEdit("0")
        self.lineEditHymn2.setValidator(validNumber)
        self.lineEditHymn2.setClearButtonEnabled(True)
        self.lineEditHymn2.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn2.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn2.editingFinished.connect(self.UpdateHymnNumber2)

        self.textFirstLine2 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine2.setReadOnly(True)
        self.textFirstLine2.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn2.currentIndexChanged.connect(self.BookUpdate2)

        layout2.addWidget(numberLabel2)
        layout2.addWidget(self.comboHymn2)
        layout2.addWidget(self.lineEditHymn2)
        layout2.addWidget(self.textFirstLine2)
        layoutLeft.addLayout(layout2)

        numberLabel3 = QLabel("3.")
        numberLabel3.setStyleSheet(cfg.colStr)
        numberLabel3.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel3.setMinimumWidth(MIN_NUMBER_WIDTH)
        
        self.comboHymn3 = QComboBox()
        self.comboHymn3.setEditable(False)
        self.comboHymn3.addItems(cfg.bookCodeList)
        
        self.lineEditHymn3 = QLineEdit("0")
        self.lineEditHymn3.setValidator(validNumber)
        self.lineEditHymn3.setClearButtonEnabled(True)
        self.lineEditHymn3.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn3.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn3.editingFinished.connect(self.UpdateHymnNumber3)

        self.textFirstLine3 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine3.setReadOnly(True)
        self.textFirstLine3.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn3.currentIndexChanged.connect(self.BookUpdate3)

        layout3.addWidget(numberLabel3)
        layout3.addWidget(self.comboHymn3)
        layout3.addWidget(self.lineEditHymn3)
        layout3.addWidget(self.textFirstLine3)
        layoutLeft.addLayout(layout3)

        numberLabel4 = QLabel("4.")
        numberLabel4.setStyleSheet(cfg.colStr)
        numberLabel4.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel4.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn4 = QComboBox()
        self.comboHymn4.setEditable(False)
        self.comboHymn4.addItems(cfg.bookCodeList)
        
        self.lineEditHymn4 = QLineEdit("0")
        self.lineEditHymn4.setValidator(validNumber)
        self.lineEditHymn4.setClearButtonEnabled(True)
        self.lineEditHymn4.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn4.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn4.editingFinished.connect(self.UpdateHymnNumber4)

        self.textFirstLine4 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine4.setReadOnly(True)
        self.textFirstLine4.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn4.currentIndexChanged.connect(self.BookUpdate4)

        layout4.addWidget(numberLabel4)
        layout4.addWidget(self.comboHymn4)
        layout4.addWidget(self.lineEditHymn4)
        layout4.addWidget(self.textFirstLine4)

        layoutLeft.addLayout(layout4)

        numberLabel5 = QLabel("5.")
        numberLabel5.setStyleSheet(cfg.colStr)
        numberLabel5.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel5.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn5 = QComboBox()
        self.comboHymn5.setEditable(False)
        self.comboHymn5.addItems(cfg.bookCodeList)
        
        self.lineEditHymn5 = QLineEdit("0")
        self.lineEditHymn5.setValidator(validNumber)
        self.lineEditHymn5.setClearButtonEnabled(True)
        self.lineEditHymn5.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn5.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn5.editingFinished.connect(self.UpdateHymnNumber5)

        self.textFirstLine5 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine5.setReadOnly(True)
        self.textFirstLine5.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn5.currentIndexChanged.connect(self.BookUpdate5)

        layout5.addWidget(numberLabel5)
        layout5.addWidget(self.comboHymn5)
        layout5.addWidget(self.lineEditHymn5)
        layout5.addWidget(self.textFirstLine5)
        
        layoutLeft.addLayout(layout5)

        numberLabel6 = QLabel("6.")
        numberLabel6.setStyleSheet(cfg.colStr)
        numberLabel6.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel6.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn6 = QComboBox()
        self.comboHymn6.setEditable(False)
        self.comboHymn6.addItems(cfg.bookCodeList)

        self.lineEditHymn6 = QLineEdit("0")
        self.lineEditHymn6.setValidator(validNumber)
        self.lineEditHymn6.setClearButtonEnabled(True)
        self.lineEditHymn6.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn6.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn6.editingFinished.connect(self.UpdateHymnNumber6)

        self.textFirstLine6 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine6.setReadOnly(True)
        self.textFirstLine6.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn6.currentIndexChanged.connect(self.BookUpdate6)

        layout6.addWidget(numberLabel6)
        layout6.addWidget(self.comboHymn6)
        layout6.addWidget(self.lineEditHymn6)
        layout6.addWidget(self.textFirstLine6)
        
        layoutLeft.addLayout(layout6)

        numberLabel7 = QLabel("7.")
        numberLabel7.setStyleSheet(cfg.colStr)
        numberLabel7.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        numberLabel7.setMinimumWidth(MIN_NUMBER_WIDTH)

        self.comboHymn7 = QComboBox()
        self.comboHymn7.setEditable(False)
        self.comboHymn7.addItems(cfg.bookCodeList)

        self.lineEditHymn7 = QLineEdit("0")
        self.lineEditHymn7.setValidator(validNumber)
        self.lineEditHymn7.setClearButtonEnabled(True)
        self.lineEditHymn7.setMinimumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn7.setMaximumWidth(EXACT_NUMBER_WIDTH)
        self.lineEditHymn7.editingFinished.connect(self.UpdateHymnNumber7)

        self.textFirstLine7 = QLineEdit("Set hymn book and number to check first line.")
        self.textFirstLine7.setReadOnly(True)
        self.textFirstLine7.setMinimumWidth(MIN_LINE_WIDTH)

        self.comboHymn7.currentIndexChanged.connect(self.BookUpdate7)

        layout7.addWidget(numberLabel7)
        layout7.addWidget(self.comboHymn7)
        layout7.addWidget(self.lineEditHymn7)
        layout7.addWidget(self.textFirstLine7)

        layoutLeft.addLayout(layout7)

        layoutThemeChooser = QVBoxLayout()
        dummyHeader1 = QLabel(' ')
        dummyHeader1.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE/4)) # to create a small linespace
        themeHeader = QLabel('Select theme')
        themeHeader.setStyleSheet(cfg.colStr)
        themeHeader.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))

        self.comboTheme = QComboBox()  # This one not editable. 
        self.comboTheme.addItems([
                                    'White sans serif typeface for a dark (red) background.',
                                    'White serif typeface for a dark (red) background.',
                                    'Black sans serif typeface for a light background.', 
                                    'Black serif typeface for a light background.' 
                                    ])
        self.comboTheme.currentIndexChanged.connect(self.ThemeIndexUpdate)
        layoutThemeChooser.addWidget(dummyHeader1)
        layoutThemeChooser.addWidget(themeHeader)
        layoutThemeChooser.addWidget(self.comboTheme)
        layoutLeft.addLayout(layoutThemeChooser) # add into layoutRight

        layoutButton = QVBoxLayout()
        dummyHeader2 = QLabel(' ')
        dummyHeader2.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE/4)) # to create a small linespace
        self.buttonMake = QPushButton('Make presentation')
        self.buttonMake.setStyleSheet(cfg.colStr)
        self.buttonMake.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
        self.buttonMake.clicked.connect(self.writePresentation) # darkOrLight is set by ThemeIndexUpdate

        layoutLeft.addWidget(dummyHeader2)
        layoutLeft.addWidget(self.buttonMake)
        
 
        layoutOverall = QHBoxLayout()
        layoutOverall.addLayout(layoutLeft)

        # instantiate the whole box
        
        widget = QWidget()
        widget.setLayout(layoutOverall)
        self.setCentralWidget(widget)
        
        # end __init__()

    def Reload(self):
        """
        Reload the hymnbook - e.g. after correcting or adding a hymn file mid session. 
        """
        cfg.hymnBook.clear()
        cfg.hymnBook, cfg.bookCodeList, message = cfg.load_hymns()
        dlg = QMessageBox(self)
        dlg.setWindowTitle("H2P starting")
        dlg.setFont(QFont(cfg.families[0],14))
        dlg.setText(message)
        button = dlg.exec()

    def showLog(self):
        """
        Show the log file in the webbrowser.
        """
        fullPathToLog = str(Path(cfg.pwd, 'H2P-private-data', 'hymnBook_reading_log.html'))

        webbrowser.open_new(r'file://'+fullPathToLog)

        
    class BookChoiceDialog(QDialog):
        
        def __init__(self):
            super().__init__()

            self.setWindowTitle("Choose book source.")

            layout = QVBoxLayout()
            message = QLabel("Which book source do you want to index?")
            message.setStyleSheet(cfg.colStr)
            message.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            self.combo = QComboBox()
            self.combo.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            self.combo.addItems(cfg.bookCodeList)
            lapout4buttons = QVBoxLayout()
            button1 = QPushButton("Index by number")
            button1.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            button2 = QPushButton("Index by first line")
            button2.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            button3 = QPushButton("Cancel")
            button3.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            
            button1.clicked.connect(self.bookCodeIndexNo)
            button2.clicked.connect(self.bookCodeIndexFL)
            button3.clicked.connect(self.justReturn)
            lapout4buttons.addWidget(button1)
            lapout4buttons.addWidget(button2)
            lapout4buttons.addWidget(button3)
            layout.addWidget(message)
            layout.addWidget(self.combo)
            layout.addLayout(lapout4buttons)
            self.setLayout(layout)

        def bookCodeIndexNo(self):
            """
            Get choice having pressed Ok.
            Produce index by number for the chosen book.
            Display index in webbrowser.
            """
            cfg.bookCodeChoice = self.combo.currentText()
            self.destroy()
            
            # compile a dist of {hymn_number: firstLine}
            tempIndx = {} 
            for tag in cfg.hymnBook:
                [num, book] = cfg.get_num_from_tag(tag)
                if book == cfg.bookCodeChoice:
                    tempIndx.update({num: cfg.hymnBook[tag].firstLine})
            # sort tempIndx
            sortedIndx = dict(sorted(tempIndx.items()))
            # write it to a fileSourceSerifPro-Regular.otf
            fullPathToIndx = str(Path(cfg.pwd, 'H2P-private-data', book + '-index.html'))
            fh = open(fullPathToIndx,'wt')
            fh.write("<b>Index for " + cfg.bookCodeChoice + "</b><br/>\n")
            fh.write("============<br/>\n<br/>\n")
            for key in sortedIndx: # key is an int in this dict
                fh.write("<pre>" + str(key).ljust(5) + sortedIndx[key] + "\n</pre>")  # <br/>
            fh.close()
            # display
            webbrowser.open_new(r'file://'+fullPathToIndx)
            return

        def bookCodeIndexFL(self):
            """
            Get choice having pressed Ok.
            Produce index by firstLine for the chosen book.
            Display index in webbrowser.
            NOT DONE YET - STILL SAME AS BY NO.
            """
            cfg.bookCodeChoice = self.combo.currentText()
            self.destroy()
            
            # compile a dist of {hymn_number: firstLine}
            tempIndx = {} 
            for tag in cfg.hymnBook:
                [num, book] = cfg.get_num_from_tag(tag)
                if book == cfg.bookCodeChoice:
                    fL = cfg.hymnBook[tag].firstLine
                    if chr(39) in fL: # get rid of ' or " in first line. spaces already stripped
                        fL = fL.replace(chr(39),"")
                    if chr(34) in fL:
                        fL = fL.replace(chr(34),"")

                    tempIndx.update({num: fL})
            # sort indx
            sortedIndx = dict(sorted(tempIndx.items(), 
                          key=lambda item: item[1]))
            # write it to a fileSourceSerifPro-Regular.otf
            fullPathToIndx = str(Path(cfg.pwd, 'H2P-private-data', book + '-index.html'))
            fh = open(fullPathToIndx,'wt')
            fh.write("<b>Index for " + cfg.bookCodeChoice + "</b><br/>\n") # Use both "\n" and HTML <\br> for good measure.
            fh.write("============<br/>\n<br/>\n")                 # Incase we want ot look at the file in an editor.
            for key in sortedIndx: # key is an int in this dict
                fh.write("<pre>" + sortedIndx[key].ljust(55) + str(key).ljust(5) + "</pre/>\n") # with pre there is no need for <br>
            fh.close()
            # display
            webbrowser.open_new(r'file://'+fullPathToIndx)
            return
            
        def justReturn(self):
            self.destroy()
            return

    
    def MakeIndex(self):
        """
        #Choose one of the source books from cfg.bookCodeList
        #and make a sorted list of all hymns in the cfg.hymnBook
        #with that booCode.
        """
        dlg = self.BookChoiceDialog()
        dlg.exec()


    class CCLIDialog(QDialog):
        
        def __init__(self):
            super().__init__()

            self.setWindowTitle("Choose CCLI number.")

            layout = QVBoxLayout()
            message = QLabel("Select church \n" + \
             "or add new church at the bottom of the list.\n" + \
             "Can delete to remove entry. \n" + \
             "(Use no digit characters in church name \n" + \
             "and only digit characters in number).\n")
            message.setStyleSheet(cfg.colStr)
            message.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            self.combo = QComboBox()
            self.combo.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            self.combo.setEditable(True)
            self.combo.addItems(cfg.CCLIList)
            completer = QCompleter(cfg.CCLIList)
            self.combo.setCompleter(completer)
            lapout4buttons = QHBoxLayout()
            button1 = QPushButton("Set/update")
            button1.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            button2 = QPushButton("Cancel")
            button2.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            button1.clicked.connect(self.setCCLInumber)
            button2.clicked.connect(self.justReturn)
            lapout4buttons.addWidget(button1)
            lapout4buttons.addWidget(button2)
            layout.addWidget(message)
            layout.addWidget(self.combo)
            layout.addLayout(lapout4buttons)
            self.setLayout(layout)

        def setCCLInumber(self):
            """
            belongs to setCCLI dialogue
            set the CCLI number in cfg.CCLI
            """
            nameAndCCLI = self.combo.currentText().strip()
            indx = self.combo.currentIndex()

            # update entry in CCLIList to current text (it may or may not have been editted
            # set cfg.ccliNumber
            # if it is now empty delete this entry from the list for next time.
            # if it is the last entry in the list and has been changed add another 'Add..' line
            # save the new list to file
            # for clarity - issue meesageBox saying what the CCLI number is being set to.
            # destroy the dialog
                
            cfg.CCLIList[indx] = nameAndCCLI
            cfg.ccliNumber = ''.join(char for char in nameAndCCLI if char.isdigit())
            if cfg.ccliNumber == '': # can find no number in whatever the (un)editted text is
                cfg.ccliNumber = '0' # so can detect not to display
                
            if ('' in cfg.CCLIList):
                cfg.CCLIList.remove('') # remove any deleted entries for next time
                
            if cfg.CCLIList[-1] != 'Add a new church name and CCLI No. here.':
                    cfg.CCLIList.append('Add a new church name and CCLI No. here.') # for next time
                
            # save the new list
            
            CCLIFile = open(Path('.', 'H2P-private-data', 'ccli_number.txt'),"w", encoding = 'utf-8')
            for iline in range(0,len(cfg.CCLIList)):
                CCLIFile.write(cfg.CCLIList[iline] + "\n")
            CCLIFile.close()
            mdlg = QMessageBox(self)
            mdlg.setWindowTitle("CCLI number set")
            mdlg.setText("The CCLI number has been set to " + cfg.ccliNumber)
            mdlg.setFont(QFont(cfg.families[0],INTERFACE_FONT_SIZE))
            button = mdlg.exec()
            self.destroy()
           

        def justReturn(self):
            self.destroy()
            return
            

    def setCCLI(self):
        """
        Set the CCLI licence to be displayed at the end of all relevant
        hymns in the presentation.
        """
        dlg = self.CCLIDialog()
        dlg.exec()

    def recordBlackBetween(self):
        """
        Record the state of the checkable menu item
        to make the choice of cfg.black_between
        persist between sessions.
        """
        cfg.black_between = self.menuActionBlack.isChecked()
        blackFile = open(Path(cfg.pwd, 'H2P-private-data', 'black_between.txt'),"w", encoding = 'utf-8')
        if cfg.black_between:
            blackFile.write("Black divider\n")
        else:
            blackFile.write("Editable divider\n")
        blackFile.close()
        return

    def Help(self):
        """
        Open the software manual 
        """
        # have already imported Path, webbrowser, os at top of this source file
        
        fullPathToHelp = str(Path(cfg.pwd, 'H2P-private-data', 'H2P-help-txt.pdf'))

        webbrowser.open_new(r'file://'+fullPathToHelp)
        return

    def About(self):
        """
        Show author info 
        """
        dlg = QMessageBox(self)
        dlg.setWindowTitle("About H2P")
        dlg.setFont(QFont(cfg.families[0],12))
        dlg.setText(
        "H2P  Copyright (C) 2025  Chris Jones, vers " + cfg.vers + "\n" + \
        "This program comes with ABSOLUTELY NO WARRANTY. \n" + \
        "This is free software, and you are welcome to \n" + \
        "redistribute it under certain conditions.\n" + \
        "For details see the GPL Licence version 3 \n" + \
        "included in the set of files making up this software.\n\n" + \
        "Additionally, this software should only be used by\n" + \
        "churches that have copyright access (i.e. have bought\n" + \
        "the included hymn books) and have a current CCLI licence\n" +\
        "that confers the right to project.\n\n" + \
        "The software is written in python (vers 3).\n" + \
        "It uses the github.com/scanny/python-pptx module\n" + \
        "and the community version of pyside6. ")
        
        button = dlg.exec()
        return

    def Exit(self):
        """
        Close the whole application in a tidy way.
        Note, can do this anyway with the x button.
        """
        
        sys.exit()


    def BookUpdate1(self, indx):
        """
        Update which book code is selcted for hymn 1.
        """

        cfg.hymnBCode[0] = cfg.bookCodeList[indx]
        s = self.lineEditHymn1.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine1.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[0],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine1.setText(text)
                cfg.tagList[0] = tag
            else:
                self.textFirstLine1.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine1.setText("Set hymn book and number to check first line.")
            cfg.tagList[0] = tag # even though we have woked out the hymn number is 0
            


    def BookUpdate2(self, indx):
        """
        Update which book code is selcted for hymn 2.
        """

        cfg.hymnBCode[1] = cfg.bookCodeList[indx]
        s = self.lineEditHymn2.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine2.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[1],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine2.setText(text)
                cfg.tagList[1] = tag
            else:
                self.textFirstLine2.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine2.setText("Set hymn book and number to check first line.")
            cfg.tagList[1] = tag # even though we have woked out the hymn number is 0
            

    def BookUpdate3(self, indx):
        """
        Update which book code is selcted for hymn 3.
        """

        cfg.hymnBCode[2] = cfg.bookCodeList[indx]
        s = self.lineEditHymn3.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine3.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[2],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine3.setText(text)
                cfg.tagList[2] = tag
            else:
                self.textFirstLine3.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine3.setText("Set hymn book and number to check first line.")
            cfg.tagList[2] = tag # even though we have woked out the hymn number is 0

            
    def BookUpdate4(self, indx):
        """
        Update which book code is selcted for hymn 4.
        """

        cfg.hymnBCode[3] = cfg.bookCodeList[indx]
        s = self.lineEditHymn4.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine4.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[3],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine4.setText(text)
                cfg.tagList[3] = tag
            else:
                self.textFirstLine4.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine4.setText("Set hymn book and number to check first line.")
            cfg.tagList[3] = tag # even though we have woked out the hymn number is 0


    def BookUpdate5(self, indx):
        """ 
        Update which book code is selcted for hymn 5.
        """
        
        cfg.hymnBCode[4] = cfg.bookCodeList[indx]
        s = self.lineEditHymn5.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine5.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[4],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine5.setText(text)
                cfg.tagList[4] = tag
            else:
                self.textFirstLine5.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine5.setText("Set hymn book and number to check first line.")
            cfg.tagList[4] = tag # even though we have woked out the hymn number is 0


    def BookUpdate6(self, indx):
        """
        Update which book code is selcted for hymn 6.
        """
        cfg.hymnBCode[5] = cfg.bookCodeList[indx]
        s = self.lineEditHymn6.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine6.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[5],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine6.setText(text)
                cfg.tagList[5] = tag
            else:
                self.textFirstLine6.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine6.setText("Set hymn book and number to check first line.")
            cfg.tagList[5] = tag # even though we have woked out the hymn number is 0

    def BookUpdate7(self, indx):
        """
        Update which book code is selcted for hymn 7.
        """
        
        cfg.hymnBCode[6] = cfg.bookCodeList[indx]
        s = self.lineEditHymn7.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine7.setText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[6],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine7.setText(text)
                cfg.tagList[6] = tag
            else:
                self.textFirstLine6.setText("Hymn " + tag + " not in database.")
        else:
            self.textFirstLine7.setText("Set hymn book and number to check first line.")
            cfg.tagList[6] = tag # even though we have woked out the hymn number is 0
        
        
    def UpdateHymnNumber1(self):
        """
        Update the hymn number chosen for hymn.
        """

        s = self.lineEditHymn1.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine1.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[0],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine1.setText(text)
                cfg.tagList[0] = tag
            else:
                self.textFirstLine1.setText("Hymn " + tag + " not in database.")
                cfg.tagList[0] = tag
        else:
            self.textFirstLine1.setText("Hymn number not set")
            cfg.tagList[0] = tag # even though we have woked out the hymn number is 0


    def UpdateHymnNumber2(self):
        """
        Update the hymn number chosen for hymn 2.
        """

        s = self.lineEditHymn2.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine2.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[1],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine2.setText(text)
                cfg.tagList[1] = tag
            else:
                self.textFirstLine2.setText("Hymn " + tag + " not in database.")
                cfg.tagList[1] = tag
        else:
            self.textFirstLine2.setText("Hymn number not set")
            cfg.tagList[1] = tag # even though we have woked out the hymn number is 0
            

    def UpdateHymnNumber3(self):
        """
        Update the hymn number chosen for hymn 3.
        """

        s = self.lineEditHymn3.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine3.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[2],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine3.setText(text)
                cfg.tagList[2] = tag
            else:
                self.textFirstLine3.setText("Hymn " + tag + " not in database.")
                cfg.tagList[2] = tag
        else:
            self.textFirstLine3.setText("Hymn number not set")
            cfg.tagList[2] = tag # even though we have woked out the hymn number is 0

            
    def UpdateHymnNumber4(self):
        """
        Update the hymn number chosen for hymn 4.
        """

        s = self.lineEditHymn4.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine4.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[3],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine4.setText(text)
                cfg.tagList[3] = tag
            else:
                self.textFirstLine4.setText("Hymn " + tag + " not in database.")
                cfg.tagList[3] = tag
        else:
            self.textFirstLine4.setText("Hymn number not set")
            cfg.tagList[3] = tag # even though we have woked out the hymn number is 0


    def UpdateHymnNumber5(self):
        """
        Update the hymn number chosen for hymn 5.
        """

        s = self.lineEditHymn5.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine5.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[4],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine5.setText(text)
                cfg.tagList[4] = tag
            else:
                self.textFirstLine5.setText("Hymn " + tag + " not in database.")
                cfg.tagList[4] = tag
        else:
            self.textFirstLine5.setText("Hymn number not set")
            cfg.tagList[4] = tag # even though we have woked out the hymn number is 0


    def UpdateHymnNumber6(self):
        """
        Update the hymn number chosen for hymn 6.
        """

        s = self.lineEditHymn6.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine6.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[5],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine6.setText(text)
                cfg.tagList[5] = tag
            else:
                self.textFirstLine6.setText("Hymn " + tag + " not in database.")
                cfg.tagList[5] = tag
        else:
            self.textFirstLine6.setText("Hymn number not set")
            cfg.tagList[5] = tag # even though we have woked out the hymn number is 0


    def UpdateHymnNumber7(self):
        """
        Update the hymn number chosen for hymn 7.
        """

        s = self.lineEditHymn7.text().strip()
        if len(s) > 0:
            s_num = int(s)
        else:
            s_num = 0
            self.textFirstLine7.setPlaceholderText("Set hymn book and number to check first line.")

        tag, bookCode, numberBoucedBack = cfg.make_tag_from_num(cfg.hymnBCode[6],s_num)
        if s_num > 0:
            if tag in cfg.hymnBook:
                text = cfg.hymnBook[tag].firstLine
                self.textFirstLine7.setText(text)
                cfg.tagList[6] = tag
            else:
                self.textFirstLine7.setText("Hymn " + tag + " not in database.")
                cfg.tagList[6] = tag
        else:
            self.textFirstLine7.setText("Hymn number not set")
            cfg.tagList[6] = tag # even though we have woked out the hymn number is 0


    def ThemeIndexUpdate(self, indx):
        """
        Update which basic theme is chosen in the comboTheme widget.
        This is only for keeping comboBox sync'd.
        The comboBox is read again when we open_presentation().
        """

        cfg.darkOrLight = indx # takes values 0, 1 ,2 ,3
            
        if indx == 0 or indx == 2:
            cfg.fontName = 'Andika' 
        else:
            cfg.fontName = 'Source Serif Pro'


    def writePresentation(self):
        """
        Connected as the callback to the makeButton.
        Make the presentation and save it on the desktop.
        """

        from pptx import Presentation
        from pptx.util import Inches, Pt, Cm
        from pptx.enum.text import PP_ALIGN


        def open_presentation():
    
            """
            Local fn to writePresentation
            
            Get one of two (only) specially prepared presentation filels with only a title page and a
            blank slide. These are prepared using freeoffice in which a standard title and content page
            with placeholders can be turned into a blank page by deleting the placeholder.
            Do not loose the two base pptx files.

            see https://python-pptx-fix.readthedocs.io/en/stable/index.html#
            
            """

            # globals shared by open_presentation(), inject_verse(), inject_hymn_number(), inject_attribution, inject_blank_slide
            
            global prs, verse_slide_layout, between_slide_layout, black_slide_layout, color_num

            # get the presentation indx again. It may never have been retrieved if combo box not used.

            cfg.black_between = self.menuActionBlack.isChecked()
            
            indx = self.comboTheme.currentIndex()

            cfg.darkOrLight = indx # takes values 0, 1, 2 ,3
                
            if indx == 0 or indx == 1:
                cfg.fontName = 'Andika' 
            else:
                cfg.fontName = 'Source Serif Pro'

            match cfg.darkOrLight:
                case 0 :
                    prs = Presentation(Path(cfg.pwd, 'H2P-private-data', 'dark_base_sans.pptx'))
                    cfg.fontName = 'Andika' 
                    color_num = 1 # this is a number in the starting presetations colour scheme. 
                case 1 :
                    prs = Presentation(Path(cfg.pwd, 'H2P-private-data', 'dark_base_serif.pptx'))
                    color_num = 1
                    cfg.fontName = 'Source Serif Pro'
                case 2 :
                    prs = Presentation(Path(cfg.pwd, 'H2P-private-data', 'light_base_sans.pptx'))
                    color_num = 1
                    cfg.fontName = 'Andika' 
                case 3 :
                    prs = Presentation(Path(cfg.pwd, 'H2P-private-data', 'light_base_serif.pptx'))
                    color_num = 1
                    cfg.fontName = 'Source Serif Pro'
                case _ : # should never get here as values set by comboBox, but default
                    prs = Presentation(Path(cfg.pwd, 'H2P-private-data', 'dark_base_sans.pptx'))
                    color_num = 1
                    cfg.fontName = 'Andika' 
                
            # These 4 presentations have already been set up in Impress/Freeoffice to be docile
            # and 16 by 9 format.
            # There is already one slide in the presentation - a title slide, leave it there
            # These presentations have only one (title) slide.
            # Their Master, however, has 3 layouts. I read these as verse_slide_layout,
            # between_slide_layout and black_slide_layout. The latter has a placeholder.
            # I have found FreeOffice to be the best s/w
            # for gaining access to and modifying layouts and placeholders within them.

            slide0 = prs.slides[0]
                
            verse_slide_layout = prs.slide_layouts[0] # this is the case only in the special base files
            between_slide_layout = prs.slide_layouts[1]
            black_slide_layout = prs.slide_layouts[2]
            
            return


        def inject_verse(verse, fontSize, isChorus = False):
            """
            Note that inject verse generates the new slide so
            adding attribute box or hymn number box must follow this and pick up the current slide.

            This fn used for chorus as well but chorus set to italic, see below.
            """
            def prepLine(line: str):
                """
                Prepare a line of text and switch on bold, italics or forced line-space.
                This involves examinging the first character of the line.

                If the line has only ~ or # or * and no other non-white-space character
                it it possible that pline = '', i.e. empty
                """
                bold = False
                italic = False
                forceLineSpace = False
                pline = line.strip() # in case there is a space before the * etc.

                if pline[0] == '*':
                    bold = True
                    pline = line[1:len(line)]
                    pline.strip() # incase there was a space after * etc.
                elif line[0] == '#':
                    italic = True
                    pline = line[1:len(line)]
                    pline.strip()
                elif line[0] == '~':
                    pline = line[1:len(line)]
                    pline.strip()
                    pline = chr(9) + pline
                elif line[0] == '$':
                    forceLineSpace = True
                    pline = ' '
                if isChorus:
                    italic = True
                    
                return pline, bold, italic, forceLineSpace

            global prs, verse_slide_layout, between_slide_layout, color_num


            # dimensions of a verse box (local variables)
            left = 1.0 * Cm(1)
            top = 1.8 * Cm(1)
            width = 26.5 * Cm(1) # this 0.5 cm > I used to set pt sizes for maxLineLength
            height = 12.4 * Cm(1) # also small leeway here

            # generate a new slide
            titleSLide = prs.slides[0]
            
            currentSlide = prs.slides.add_slide(verse_slide_layout)
             
            txBox1 = currentSlide.shapes.add_textbox(left, top, width, height)
            tf1 = txBox1.text_frame
            p = []
            run = []
            font = []

            for iline in  range(0,len(verse)):
                pline, bold, italic, forceLineSpace = prepLine(verse[iline])
                if iline == 0: 
                    p.append(tf1.paragraphs[0]) # p[0]
                else:
                    p.append(tf1.add_paragraph())
                if len(pline) > 0: # because of poss. pline = '', see prepLine()
                    p[iline].alignment = PP_ALIGN.LEFT
                    run.append(p[iline].add_run())
                    run[iline].text =  pline
                    font.append(run[iline].font)
                    font[iline].name = cfg.fontName
                    if forceLineSpace:
                        font[iline].size = Pt(14)
                    else:
                        font[iline].size = Pt(fontSize)
                    font[iline].bold = bold
                    font[iline].italic = italic
                    font[iline].color.theme_color = color_num # 1 is black, 2 is white -- that is all
 
            return currentSlide

        def inject_attribution(currentSlide, attribution, rcode):

            global prs, verse_slide_layout, between_slide_layout, black_slide_layout, color_num

            # dimensions of a verse box
            left = 1.0 * Cm(1)
            top = 14.5 * Cm(1)
            width = 26 * Cm(1)
            height = 1 * Cm(1)
            fontSize = 18
            txBox1 = currentSlide.shapes.add_textbox(left, top, width, height)
            tf1 = txBox1.text_frame
            p1 = tf1.paragraphs[0] # there is already a paragraph there when the frame when it is created.
            p1.alignment = PP_ALIGN.RIGHT
            run1 = p1.add_run()
            if hymn.rCode.upper() == "CC0":
                cp = "" # just don't mention it approach here, chr(127342) is unicode for copyright with slash through
            elif hymn.rCode.upper() == "CCLI":
                cp = "  CCLI: " + cfg.ccliNumber
            else: # if <..> is absent in hymn.txt file, rCode will have been set to ""
                cp =''
            if int(cfg.ccliNumber) != 0:    
                run1.text = attribution + cp
            else:
                run1.text = attribution
            font1 = run1.font
            font1.name = cfg.fontName
            font1.size = Pt(fontSize)
            font1.color.theme_color = color_num # 1 is black, 2 is white -- that is all
            font1.italic = False
            return

        def inject_hymn_number(currentSlide, tag, iverse, chorus = False):

            global prs, verse_slide_layout, between_slide_layout, black_slide_layout, color_num

            # dimensions of a hymn number box box
            left = 16.0 * Cm(1)
            top = 0.3 * Cm(1)
            width = 11 * Cm(1)
            height = 1.2 * Cm(1)
            fontSize = 28 # for Hymn number

            txBox1 = currentSlide.shapes.add_textbox(left, top, width, height)
            tf1 = txBox1.text_frame
            p1 = tf1.paragraphs[0]
            p1.alignment = PP_ALIGN.RIGHT
            run1 = p1.add_run()
            outOf = cfg.hymnBook[tag].topBookVerseNumber    
            if not bool(cfg.hymnBook[tag].bookVerseNumber): # forced verse numbers dict is empty
                if len(cfg.hymnBook[tag].bookVerseNumber) == 1: # just one verse
                    verseNumber = ''
                else:
                    verseNumber = ' v ' + str(iverse+1) + '/' + str(outOf)
                if chorus:
                    verseNumber = verseNumber + " c"
            elif iverse in cfg.hymnBook[tag].bookVerseNumber: # there are forced numberings
                verseNumber = cfg.hymnBook[tag].bookVerseNumber[iverse] 
            else: 
                verseNumber = ''
            if iverse == -1: # set like this for chorus first
                verseNumber = ' c'
            run1.text = tag + verseNumber
            font1 = run1.font
            font1.name = cfg.fontName
            font1.size = Pt(fontSize)
            font1.color.theme_color = color_num # 1 is black, 2 is white -- that is all
            font1.italic = True
            return


        def inject_blank_slide(black: bool):
            """.. but put a text box in it.
            """

            global prs, verse_slide_layout, between_slide_layout, black_slide_layout, color_num

            # dimensions of a text box
            left = 8.5 * Cm(1)
            top = 5.5 * Cm(1)
            width = 12 * Cm(1)
            height = 2 * Cm(1)
            fontSize = 40
            if black:
                currentSlide = prs.slides.add_slide(black_slide_layout)
            else:
                currentSlide = prs.slides.add_slide(between_slide_layout)

            return currentSlide


        # self.buttonMake.setText('Working')
        self.buttonMake.setStyleSheet("color: red") # flash red to make action more visible (windows seems to need this)
        QApplication.processEvents()
        open_presentation()

        for tag in cfg.tagList:
            if tag in cfg.hymnBook:
                hymn = cfg.hymnBook[tag]
                fontSize = hymn.fontSize
                nslides = len(hymn.verse)
                if hymn.hasChorus:
                    nslides = 2 * nslides
                slideCount = 0
                if hymn.chorusFirst:  # chorus is written first
                    slideCount += 1
                    currentSlide = inject_verse(hymn.chorus, hymn.fontSize, isChorus = True)
                    inject_hymn_number(currentSlide, tag, -1)
                    
                for iverse in range(0,len(hymn.verse)):
                    slideCount += 1
                    currentSlide = inject_verse(hymn.verse[iverse], hymn.fontSize)
                    inject_hymn_number(currentSlide, tag, iverse) 
                    if hymn.hasChorus:
                        slideCount += 1
                        currentSlide = inject_verse(hymn.chorus, hymn.fontSize, isChorus = True)
                        inject_hymn_number(currentSlide, tag, iverse, chorus = True)
                    if iverse == len(hymn.verse) - 1: # last slide of hymn
                        inject_attribution(currentSlide, hymn.attribution, hymn.rCode.upper())
                        
                inject_blank_slide(cfg.black_between)
            elif cfg.get_num_from_tag(tag)[0] > 0:
                text = []
                text.append(tag + " not in database yet.") # text must mimic a verse and therefore be a list of strings
                fontSize = 42
                currentSlide = inject_verse(text, fontSize)
                inject_blank_slide(cfg.black_between)

        # it seems convenient to provide one last editable-layout slide at the end
        if cfg.black_between:
            inject_blank_slide(False)
                
        prs.save(Path(cfg.desktop,'hymn-words.pptx'))
        
        time.sleep(0.3)
        # self.buttonMake.setText("Make presentation")
        self.buttonMake.setStyleSheet(cfg.colStr) # back to normal colour
        QApplication.processEvents()

